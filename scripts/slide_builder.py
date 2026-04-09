"""
slide_builder.py
Generates a PowerPoint presentation from faculty-specified topics,
grounded in content retrieved from the vector DB.
"""

import os
import json
import asyncio
from pathlib import Path
from typing import Optional

from pydantic import BaseModel, Field
from pptx import Presentation
from pptx.util import Inches, Pt
from langchain_google_cloud_sql_pg import PostgresVectorStore, PostgresEngine
from langchain_openai import OpenAIEmbeddings
from langchain_groq import ChatGroq
from dotenv import load_dotenv

load_dotenv()

# ---------------------------------------------------------------------------
# Config (mirrors query_data.py)
# ---------------------------------------------------------------------------

PROJECT_ID  = os.getenv("GCP_PROJECT_ID")
REGION      = os.getenv("GCP_REGION")
INSTANCE    = os.getenv("CLOUD_SQL_INSTANCE")
DATABASE    = os.getenv("CLOUD_SQL_DATABASE")
DB_USER     = os.getenv("CLOUD_SQL_USER")
DB_PASS     = os.getenv("CLOUD_SQL_PASSWORD")
TABLE_NAME  = "course_embeddings"
GROQ_API_KEY = os.getenv("GROQ_API_KEY")


# ---------------------------------------------------------------------------
# Pydantic Models
# ---------------------------------------------------------------------------

class Slide(BaseModel):
    title: str = Field(..., description="Slide title")
    bullets: list[str] = Field(..., description="3-5 concise bullet points")
    speaker_notes: str = Field(default="", description="Optional speaker notes")


class PresentationModel(BaseModel):
    presentation_title: str
    course_name: str
    slides: list[Slide]


# ---------------------------------------------------------------------------
# Vector DB Retrieval
# ---------------------------------------------------------------------------

async def retrieve_chunks(topic: str, k: int = 5) -> list[dict]:
    """Retrieve relevant chunks from the vector DB for a given topic."""
    engine = await PostgresEngine.afrom_instance(
        project_id=PROJECT_ID,
        region=REGION,
        instance=INSTANCE,
        database=DATABASE,
        user=DB_USER,
        password=DB_PASS,
    )

    vector_store = await PostgresVectorStore.create(
        engine,
        embedding_service=OpenAIEmbeddings(),
        table_name=TABLE_NAME,
    )

    results = await vector_store.asimilarity_search(topic, k=k)

    return [
        {"text": doc.page_content, "metadata": doc.metadata}
        for doc in results
    ]


# ---------------------------------------------------------------------------
# LLM Setup
# ---------------------------------------------------------------------------

def get_llm() -> ChatGroq:
    return ChatGroq(
        model="llama-3.1-8b-instant",
        temperature=0.3,
        api_key=GROQ_API_KEY,
    )


# ---------------------------------------------------------------------------
# Prompt Builder
# ---------------------------------------------------------------------------

def build_prompt(topic: str, context_chunks: list[dict], slides_per_topic: int) -> str:
    context_text = "\n\n".join(
        f"[Chunk {i+1}]: {chunk['text']}"
        for i, chunk in enumerate(context_chunks)
    )

    return f"""You are an expert academic instructor preparing lecture slides.

Using ONLY the provided context below, generate exactly {slides_per_topic} slide(s) for the topic: "{topic}".

Each slide must have:
- A clear, concise title
- 3 to 5 bullet points (each under 15 words)
- Optional speaker notes expanding on the bullets

Return your response as a valid JSON array of slide objects.
Do NOT include any explanation, markdown, or text outside the JSON.

Format:
[
  {{
    "title": "Slide Title",
    "bullets": ["Bullet 1", "Bullet 2", "Bullet 3"],
    "speaker_notes": "Optional notes here."
  }}
]

Context:
{context_text}
"""


# ---------------------------------------------------------------------------
# LLM Slide Generation (per topic)
# ---------------------------------------------------------------------------

def generate_slides_for_topic(
    topic: str,
    context_chunks: list[dict],
    slides_per_topic: int,
    llm: ChatGroq,
) -> list[Slide]:
    prompt = build_prompt(topic, context_chunks, slides_per_topic)

    try:
        response = llm.invoke(prompt)
        raw = response.content.strip()

        # Strip markdown fences if present
        if raw.startswith("```"):
            raw = raw.split("```")[1]
            if raw.startswith("json"):
                raw = raw[4:]
        raw = raw.strip()

        slide_dicts = json.loads(raw)
        return [Slide(**s) for s in slide_dicts]

    except (json.JSONDecodeError, KeyError, ValueError) as e:
        print(f"[Warning] Failed to parse slides for topic '{topic}': {e}")
        return [Slide(
            title=f"{topic} (Content Unavailable)",
            bullets=["Could not generate content for this topic."],
            speaker_notes="",
        )]


# ---------------------------------------------------------------------------
# PPTX Rendering
# ---------------------------------------------------------------------------

def add_title_slide(prs: Presentation, title: str, course_name: str) -> None:
    layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = course_name


def add_content_slide(prs: Presentation, slide_data: Slide) -> None:
    layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(layout)

    slide.shapes.title.text = slide_data.title

    tf = slide.placeholders[1].text_frame
    tf.clear()

    for i, bullet in enumerate(slide_data.bullets):
        if i == 0:
            tf.paragraphs[0].text = bullet
            tf.paragraphs[0].level = 0
        else:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0

    if slide_data.speaker_notes:
        slide.notes_slide.notes_text_frame.text = slide_data.speaker_notes


def render_pptx(presentation_data: PresentationModel, output_path: str) -> str:
    prs = Presentation()
    add_title_slide(prs, presentation_data.presentation_title, presentation_data.course_name)
    for slide in presentation_data.slides:
        add_content_slide(prs, slide)
    prs.save(output_path)
    return output_path


# ---------------------------------------------------------------------------
# Main Builder Function
# ---------------------------------------------------------------------------

def build_presentation(
    topics: list[str],
    presentation_title: str,
    course_name: str,
    slides_per_topic: int = 2,
    top_k: int = 5,
    output_path: Optional[str] = None,
) -> str:
    """
    Main entry point for slide generation.

    Args:
        topics:               List of faculty-specified topics.
        presentation_title:   Title shown on the cover slide.
        course_name:          Course name shown on the cover slide.
        slides_per_topic:     How many slides to generate per topic.
        top_k:                Number of chunks to retrieve per topic from vector DB.
        output_path:          Where to save the .pptx. Defaults to outputs/<title>.pptx

    Returns:
        Path to the saved .pptx file.
    """
    if output_path is None:
        Path("outputs").mkdir(exist_ok=True)
        safe_title = presentation_title.replace(" ", "_").lower()
        output_path = f"outputs/{safe_title}.pptx"

    llm = get_llm()
    all_slides: list[Slide] = []

    for topic in topics:
        print(f"[slide_builder] Retrieving chunks for topic: '{topic}'")
        chunks = asyncio.run(retrieve_chunks(topic, k=top_k))

        if not chunks:
            print(f"[Warning] No chunks found for topic: '{topic}'. Skipping.")
            continue

        print(f"[slide_builder] Generating slides for topic: '{topic}'")
        topic_slides = generate_slides_for_topic(
            topic=topic,
            context_chunks=chunks,
            slides_per_topic=slides_per_topic,
            llm=llm,
        )
        all_slides.extend(topic_slides)

    if not all_slides:
        raise ValueError("No slides were generated. Check your topics and vector DB content.")

    pres = PresentationModel(
        presentation_title=presentation_title,
        course_name=course_name,
        slides=all_slides,
    )

    saved_path = render_pptx(pres, output_path)
    print(f"[slide_builder] Presentation saved to: {saved_path}")
    return saved_path


# ---------------------------------------------------------------------------
# Example usage
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    build_presentation(
        topics=["Neural Networks", "Backpropagation", "Overfitting"],
        presentation_title="Introduction to Deep Learning",
        course_name="CS 6120 - Advanced Machine Learning",
        slides_per_topic=2,
        top_k=5,
    )